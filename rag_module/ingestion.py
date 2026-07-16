import logging
import json
from pathlib import Path
from typing import List, Dict, Any
import pandas as pd
from pptx import Presentation
from llama_index.core import Document
from llama_index.readers.file import PDFReader, UnstructuredReader, DocxReader
from rag_module.service import rag_service
from rag_module.chunking import StructuredChunker, file_fingerprint
from config.settings import settings
from datetime import datetime

# OCR 相关导入
try:
  import pdfplumber
  PDFPLUMBER_AVAILABLE = True
except ImportError:
  PDFPLUMBER_AVAILABLE = False

try:
  from pdf2image import convert_from_path
  import pytesseract
  # 配置 Tesseract 路径（仅 Windows 需指定，Linux 用系统 PATH）
  import platform
  if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
  OCR_AVAILABLE = True
except ImportError:
  OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

class DataIngestor:
  def __init__(self):
    self.raw_dir = Path(settings.RAG_RAW_DIR)
    self.parsed_dir = Path(settings.RAG_PARSED_DIR)
    self.meta_dir = Path(settings.RAG_META_DIR)

    # 文件索引路径（记录已导入的文件）
    self.index_file = self.raw_dir / ".ingestion_index.json"
    self.file_index = self._load_index()

  def _load_index(self) -> Dict[str, Dict[str, Any]]:
    """加载文件索引"""
    if self.index_file.exists():
      with open(self.index_file, 'r', encoding='utf-8') as f:
        return json.load(f)
    return {}

  def _save_index(self):
    """保存文件索引"""
    with open(self.index_file, 'w', encoding='utf-8') as f:
      json.dump(self.file_index, f, indent=2, ensure_ascii=False)

  def _file_key(self, file_path: Path) -> str:
    for prefix, root in (("raw", self.raw_dir), ("meta", self.meta_dir)):
      try:
        return f"{prefix}/{file_path.relative_to(root).as_posix()}"
      except ValueError:
        continue
    raise ValueError(f"文件不在知识库目录中: {file_path}")

  def _get_file_hash(self, file_path: Path) -> str:
    return file_fingerprint(file_path)

  def _is_file_changed(self, file_path: Path) -> bool:
    """检查文件是否已导入且未修改"""
    file_key = self._file_key(file_path)

    if file_key not in self.file_index:
      return True  # 新文件

    # 检查哈希值
    current_hash = self._get_file_hash(file_path)
    return current_hash != self.file_index[file_key].get('hash')

  def _record_file(self, file_path: Path, doc_count: int, chunk_count: int, version: str):
    """记录已导入的文件"""
    file_key = self._file_key(file_path)
    self.file_index[file_key] = {
      'hash': version,
      'imported_at': datetime.now().isoformat(),
      'doc_count': doc_count,
      'chunk_count': chunk_count,
      'file_name': file_path.name
    }
    self._save_index()

  def _parse_file(self, file_path: Path) -> List[Document]:
    suffix = file_path.suffix.lower()
    if suffix == '.pdf':
      return self._parse_pdf(file_path)
    if suffix == '.pptx':
      return self._parse_pptx(file_path)
    if suffix in ['.xlsx', '.xls']:
      return self._parse_excel(file_path)
    if suffix == '.csv':
      return self._parse_csv(file_path)
    if suffix == '.json':
      return self._parse_json(file_path)
    if suffix in ['.txt', '.md']:
      return self._parse_txt(file_path)
    if suffix in ['.doc', '.docx']:
      return self._parse_docx(file_path)
    logger.warning(f"❓ 不支持的文件格式: {file_path.name}")
    return []

  def ingest_all_formats(self, force_reimport: bool = False):
    """统一入口：自动识别并导入所有格式的文件

    Args:
      force_reimport: 是否强制重新导入所有文件（忽略索引）
    """
    logger.info(f"开始扫描并导入 {self.raw_dir} 与 {self.meta_dir} 下的文档...")

    if force_reimport:
      logger.warning("⚠️  强制重新导入模式：将重新导入所有文件")

    if not self.raw_dir.exists() and not self.meta_dir.exists():
      logger.error("原始文档目录和实验元数据目录都不存在")
      return

    rag_service.initialize()
    if rag_service.collection_count() == 0 and self.file_index:
      logger.warning("向量集合为空，忽略旧的入库状态并重新扫描全部文档")
      self.file_index = {}
    chunker = StructuredChunker(
      parent_chunk_size=settings.RAG_PARENT_CHUNK_SIZE,
      child_chunk_size=settings.RAG_CHILD_CHUNK_SIZE,
      child_chunk_overlap=settings.RAG_CHILD_CHUNK_OVERLAP,
    )

    total_docs = 0
    total_chunks = 0
    skipped_files = 0
    new_files = 0
    updated_files = 0

    seen_file_keys = set()
    roots = [root for root in (self.raw_dir, self.meta_dir) if root.exists()]
    for file_path in (path for root in roots for path in root.rglob('*')):
      if file_path.is_file() and not file_path.name.startswith('.'):
        file_key = self._file_key(file_path)
        seen_file_keys.add(file_key)
        # 检查是否需要导入
        if not force_reimport and not self._is_file_changed(file_path):
          skipped_files += 1
          logger.debug(f"⏭️  跳过未修改文件: {file_path.name}")
          continue

        # 判断是新文件还是更新文件
        is_new = file_key not in self.file_index

        try:
          docs = self._parse_file(file_path)
          if docs:
            version = self._get_file_hash(file_path)
            chunks = chunker.chunk_documents(
              docs,
              file_key=file_key,
              ingestion_version=version,
            )
            if not chunks:
              logger.warning(f"未生成有效文本块: {file_path.name}")
              continue
            chunk_count = rag_service.upsert_chunks(chunks)
            rag_service.prune_source_versions(file_key, version)
            total_docs += len(docs)
            total_chunks += chunk_count

            # 记录文件
            self._record_file(file_path, len(docs), chunk_count, version)

            if is_new:
              new_files += 1
              logger.info(f"✅ 新增 {file_path.name}: {len(docs)} 个文档 / {chunk_count} 个子块")
            else:
              updated_files += 1
              logger.info(f"🔄 更新 {file_path.name}: {len(docs)} 个文档 / {chunk_count} 个子块")

        except Exception as e:
          logger.error(f"❌ 解析失败 {file_path.name}: {str(e)}")

    removed_keys = set(self.file_index) - seen_file_keys
    for file_key in removed_keys:
      rag_service.delete_source(file_key)
      self.file_index.pop(file_key, None)
      logger.info(f"🗑️ 已删除知识库中不存在的来源: {file_key}")

    self._save_index()

    logger.info("\n========== 导入完成 ==========")
    logger.info(f"✅ 新增文件: {new_files} 个")
    logger.info(f"🔄 更新文件: {updated_files} 个")
    logger.info(f"⏭️  跳过文件: {skipped_files} 个")
    logger.info(f"📦 共导入文档: {total_docs} 个，子块: {total_chunks} 个")
    logger.info(f"🔎 当前向量点数: {rag_service.collection_count()}")
    logger.info(f"💾 索引已保存至: {self.index_file}")

  def _parse_pdf(self, file_path: Path) -> List[Document]:
    """解析 PDF 文件（支持 OCR）"""
    documents = []

    # 方法 1: 尝试直接提取文本
    try:
      reader = PDFReader()
      documents = reader.load_data(file=str(file_path))

      # 检查是否成功提取到文本
      has_text = any(doc.text and len(doc.text.strip()) > 50 for doc in documents)

      if has_text:
        # 添加元数据
        for doc in documents:
          doc.metadata.update({
            "source_type": "pdf",
            "file_name": file_path.name,
            "file_path": str(file_path),
            "extraction_method": "text"
          })
        logger.info(f"📄 {file_path.name}: 文本提取成功")
        return documents
      else:
        logger.warning(f"⚠️  {file_path.name}: 文本提取为空，尝试 OCR...")
    except Exception as e:
      logger.warning(f"⚠️  {file_path.name}: 文本提取失败 ({e})，尝试 OCR...")

    # 方法 2: 使用 pdfplumber（更好的文本提取）
    if PDFPLUMBER_AVAILABLE:
      try:
        documents = []
        with pdfplumber.open(str(file_path)) as pdf:
          for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()

            if text and len(text.strip()) > 50:
              doc = Document(
                text=text,
                metadata={
                  "source_type": "pdf",
                  "file_name": file_path.name,
                  "file_path": str(file_path),
                  "page_label": str(page_num),
                  "extraction_method": "pdfplumber"
                }
              )
              documents.append(doc)

        if documents:
          logger.info(f"📄 {file_path.name}: pdfplumber 提取成功 ({len(documents)} 页)")
          return documents
        else:
          logger.warning(f"⚠️  {file_path.name}: pdfplumber 提取为空，尝试 OCR...")
      except Exception as e:
        logger.warning(f"⚠️  {file_path.name}: pdfplumber 失败 ({e})，尝试 OCR...")

    # 方法 3: OCR 识别（适用于扫描版 PDF）
    if OCR_AVAILABLE:
      try:
        logger.info(f"🔍 {file_path.name}: 开始 OCR 识别（可能需要几分钟）...")

        # 将 PDF 转换为图片
        images = convert_from_path(str(file_path), dpi=200)

        documents = []
        for page_num, image in enumerate(images, 1):
          # 使用 pytesseract 进行 OCR（中文+英文）
          text = pytesseract.image_to_string(image, lang='chi_sim+eng')

          if text and len(text.strip()) > 20:
            doc = Document(
              text=text,
              metadata={
                "source_type": "pdf",
                "file_name": file_path.name,
                "file_path": str(file_path),
                "page_label": str(page_num),
                "extraction_method": "ocr"
              }
            )
            documents.append(doc)

          if page_num % 10 == 0:
            logger.info(f"  OCR 进度: {page_num}/{len(images)} 页")

        if documents:
          logger.info(f"✅ {file_path.name}: OCR 识别成功 ({len(documents)} 页)")
          return documents
        else:
          logger.error(f"❌ {file_path.name}: OCR 识别失败，未提取到文本")
      except Exception as e:
        logger.error(f"❌ {file_path.name}: OCR 识别出错: {e}")
    else:
      logger.error(f"❌ {file_path.name}: OCR 库未安装，无法识别扫描版 PDF")
      logger.error("   请安装: pip install pdf2image pytesseract pillow")
      logger.error("   并安装 Tesseract-OCR: https://github.com/tesseract-ocr/tesseract")

    # 最后尝试：UnstructuredReader
    try:
      reader = UnstructuredReader()
      documents = reader.load_data(file=str(file_path))
      for doc in documents:
        doc.metadata.update({
          "source_type": "pdf",
          "file_name": file_path.name,
          "extraction_method": "unstructured"
        })
      if documents:
        logger.info(f"📄 {file_path.name}: UnstructuredReader 成功")
        return documents
    except Exception:
      logger.error(f"❌ {file_path.name}: 所有方法均失败")

    return []

  def _parse_pptx(self, file_path: Path) -> List[Document]:
    """解析 PPTX 文件"""
    prs = Presentation(str(file_path))
    documents = []

    for slide_idx, slide in enumerate(prs.slides, 1):
      text_content = []

      # 提取所有文本框内容
      for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
          text_content.append(shape.text.strip())

      if text_content:
        content = "\n".join(text_content)
        doc = Document(
          text=content,
          metadata={
            "source_type": "pptx",
            "file_name": file_path.name,
            "slide_number": slide_idx,
            "file_path": str(file_path)
          }
        )
        documents.append(doc)

    return documents

  def _parse_excel(self, file_path: Path) -> List[Document]:
    """解析 Excel 文件"""
    documents = []

    # 读取所有 sheet
    excel_file = pd.ExcelFile(str(file_path))
    for sheet_name in excel_file.sheet_names:
      df = pd.read_excel(excel_file, sheet_name=sheet_name)

      # 生成表格的语义描述
      content = f"表格名称: {sheet_name}\n"
      content += f"列名: {', '.join(df.columns.tolist())}\n"
      content += f"行数: {len(df)}\n\n"
      content += f"表格内容摘要:\n{df.head(10).to_string()}\n"

      doc = Document(
        text=content,
        metadata={
          "source_type": "excel",
          "file_name": file_path.name,
          "sheet_name": sheet_name,
          "file_path": str(file_path),
          "columns": df.columns.tolist()
        }
      )
      documents.append(doc)

    return documents

  def _parse_csv(self, file_path: Path) -> List[Document]:
    """解析 CSV 文件"""
    df = pd.read_csv(str(file_path))

    content = f"CSV 文件: {file_path.name}\n"
    content += f"列名: {', '.join(df.columns.tolist())}\n"
    content += f"行数: {len(df)}\n\n"
    content += f"数据摘要:\n{df.describe().to_string()}\n\n"
    content += f"前10行数据:\n{df.head(10).to_string()}\n"

    doc = Document(
      text=content,
      metadata={
        "source_type": "csv",
        "file_name": file_path.name,
        "file_path": str(file_path),
        "columns": df.columns.tolist()
      }
    )

    return [doc]

  def _parse_json(self, file_path: Path) -> List[Document]:
    """解析 JSON 文件"""
    with open(file_path, 'r', encoding='utf-8') as f:
      data = json.load(f)

    # 将 JSON 结构转化为可读文本
    content = f"JSON 文件: {file_path.name}\n\n"
    content += json.dumps(data, indent=2, ensure_ascii=False)

    doc = Document(
      text=content,
      metadata={
        "source_type": "json",
        "file_name": file_path.name,
        "file_path": str(file_path)
      }
    )

    return [doc]

  def _parse_txt(self, file_path: Path) -> List[Document]:
    """解析 TXT 文件"""
    try:
      # 尝试多种编码
      encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
      content = None

      for encoding in encodings:
        try:
          with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()
          break
        except UnicodeDecodeError:
          continue

      if not content:
        logger.error(f"无法读取 TXT 文件（编码问题）: {file_path.name}")
        return []

      # 如果内容太短，可能不是文档
      if len(content.strip()) < 10:
        logger.warning(f"TXT 文件内容过短，跳过: {file_path.name}")
        return []

      doc = Document(
        text=content,
        metadata={
          "source_type": "txt",
          "file_name": file_path.name,
          "file_path": str(file_path)
        }
      )

      return [doc]

    except Exception as e:
      logger.error(f"解析 TXT 文件失败: {file_path.name}, 错误: {e}")
      return []

  def _parse_docx(self, file_path: Path) -> List[Document]:
    """解析 DOC/DOCX 文件"""
    try:
      # 对于 .doc 文件，需要先转换或使用其他方法
      if file_path.suffix.lower() == '.doc':
        logger.warning(f"⚠️  .doc 格式需要手动转换为 .docx: {file_path.name}")
        logger.warning("   建议: 用 Word 打开后另存为 .docx 格式")
        return []

      # 使用 DocxReader 读取 .docx
      reader = DocxReader()
      documents = reader.load_data(file=str(file_path))

      # 添加元数据
      for doc in documents:
        doc.metadata.update({
          "source_type": "docx",
          "file_name": file_path.name,
          "file_path": str(file_path)
        })

      logger.info(f"📄 {file_path.name}: DOCX 提取成功")
      return documents

    except Exception as e:
      logger.error(f"解析 DOCX 文件失败: {file_path.name}, 错误: {e}")
      return []

ingestor = DataIngestor()
